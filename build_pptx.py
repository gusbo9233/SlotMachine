from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(0x0E, 0x10, 0x2A)
ACCENT = RGBColor(0xFF, 0xC2, 0x3C)
TEXT = RGBColor(0xF2, 0xF2, 0xF8)
SUB = RGBColor(0xB9, 0xB9, 0xD6)
CARD = RGBColor(0x1B, 0x1E, 0x40)

blank = prs.slide_layouts[6]

def add_bg(slide, color=BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.shadow.inherit = False
    return shape

def add_text(slide, left, top, width, height, text, size=18, bold=False, color=TEXT, align=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    lines = text.split('\n') if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = 'Calibri'
        if align is not None:
            p.alignment = align
    return tb

def add_card(slide, left, top, width, height, color=CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(1)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.adjustments[0] = 0.08
    return shape

def add_accent_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.18))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT

# ------- Slide 1: Title -------
s = prs.slides.add_slide(blank)
add_bg(s)
add_accent_bar(s)
add_text(s, Inches(0.8), Inches(2.2), Inches(12), Inches(1.2),
         "Casino Slots", size=64, bold=True, color=ACCENT)
add_text(s, Inches(0.8), Inches(3.4), Inches(12), Inches(0.8),
         "A React slot machine — workflow & React hooks in practice", size=26, color=TEXT)
add_text(s, Inches(0.8), Inches(6.6), Inches(12), Inches(0.4),
         "5-minute presentation  ·  React 18 + Vite", size=16, color=SUB)

# ------- Slide 2: Project overview -------
s = prs.slides.add_slide(blank)
add_bg(s)
add_accent_bar(s)
add_text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.8),
         "What is this project?", size=36, bold=True, color=ACCENT)
bullets = [
    "• Single-page React 18 app, built with Vite.",
    "• Classic 3-reel slot machine with a center payline.",
    "• 12 symbols per reel (incl. a LOGO jackpot symbol).",
    "• Player manages credits, bet size, and spin count.",
    "• Animated reels with staggered stop + win highlighting.",
    "• All logic lives in a single component: src/App.jsx.",
]
add_text(s, Inches(0.8), Inches(1.6), Inches(12), Inches(5),
         "\n".join(bullets), size=22, color=TEXT)

# ------- Slide 3: Workflow -------
s = prs.slides.add_slide(blank)
add_bg(s)
add_accent_bar(s)
add_text(s, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
         "Spin Workflow", size=36, bold=True, color=ACCENT)

steps = [
    ("1. Press Spin", "Validate credits ≥ bet, clear any active timers."),
    ("2. Deduct bet", "Pick random target index for each of the 3 reels."),
    ("3. Animate", "requestAnimationFrame cycles reel indices for ~2.2s."),
    ("4. Stagger stop", "Each reel settles 450ms after the previous one."),
    ("5. Resolve", "Read center row, detect 2x / 3x / jackpot, pay out."),
    ("6. Update UI", "Credits, last win, message, and spin counter update."),
]

y = Inches(1.4)
for i, (title, desc) in enumerate(steps):
    top = y + Inches(i * 0.9)
    add_card(s, Inches(0.6), top, Inches(12.1), Inches(0.8))
    add_text(s, Inches(0.85), top + Inches(0.08), Inches(3.2), Inches(0.7),
             title, size=20, bold=True, color=ACCENT)
    add_text(s, Inches(4.2), top + Inches(0.12), Inches(8.5), Inches(0.7),
             desc, size=18, color=TEXT)

# ------- Slide 4: useState -------
s = prs.slides.add_slide(blank)
add_bg(s)
add_accent_bar(s)
add_text(s, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
         "useState — reactive game state", size=36, bold=True, color=ACCENT)
add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "Every value the UI must re-render for lives in useState.", size=20, color=SUB)

states = [
    "centerIndices  – which symbol sits on each reel’s payline",
    "finalIndices   – locked target after a spin",
    "spinning       – are the reels currently rolling?",
    "settledReels   – per-reel stop flags for staggered finish",
    "credits / bet  – player wallet and wager",
    "message        – status text shown under the reels",
    "lastWin        – last payout info (drives highlight + glow)",
    "spinCount      – total spins this session",
]
add_text(s, Inches(0.9), Inches(2.0), Inches(12), Inches(5),
         "\n".join("• " + x for x in states), size=20, color=TEXT)

# ------- Slide 5: useEffect -------
s = prs.slides.add_slide(blank)
add_bg(s)
add_accent_bar(s)
add_text(s, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
         "useEffect — cleanup on unmount", size=36, bold=True, color=ACCENT)
add_text(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.8),
         "One effect, empty dependency array — it runs only the cleanup when App unmounts.",
         size=20, color=SUB)

add_card(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(2.4))
code = (
    "useEffect(() => {\n"
    "  return () => {\n"
    "    timersRef.current.forEach(({ type, id }) => {\n"
    "      if (type === 'frame') cancelAnimationFrame(id);\n"
    "      else clearTimeout(id);\n"
    "    });\n"
    "  };\n"
    "}, []);"
)
add_text(s, Inches(0.9), Inches(2.35), Inches(11.5), Inches(2.2),
         code, size=18, color=TEXT)

add_text(s, Inches(0.6), Inches(4.9), Inches(12), Inches(2),
         "Why it matters:\n"
         "• A spin schedules many timeouts + an rAF loop.\n"
         "• If the user navigates away mid-spin, orphan timers would\n"
         "  try to call setState on an unmounted component.\n"
         "• The cleanup cancels everything safely.",
         size=20, color=TEXT)

# ------- Slide 6: useMemo -------
s = prs.slides.add_slide(blank)
add_bg(s)
add_accent_bar(s)
add_text(s, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
         "useMemo — skip expensive rebuilds", size=36, bold=True, color=ACCENT)

add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.6),
         "Three memoised values, each with a different dependency set:", size=20, color=SUB)

items = [
    ("const reels = useMemo(buildReels, []);",
     "Builds the 3×12 symbol matrix once. Never recomputed."),
    ("const visibleGrid = useMemo(\n  () => createVisibleGrid(reels, centerIndices),\n  [centerIndices, reels]);",
     "Recomputes only when reel positions change — not on every re-render caused by bet/credits/message."),
    ("const fullStrip = useMemo(\n  () => Array.from({length: 5}, () => symbols).flat(),\n  [symbols]);   // inside <Reel>",
     "Flattens the repeated strip once per reel; animation re-renders don’t touch it."),
]
y = Inches(2.0)
for code, desc in items:
    add_card(s, Inches(0.6), y, Inches(12.1), Inches(1.55))
    add_text(s, Inches(0.85), y + Inches(0.08), Inches(6.2), Inches(1.4),
             code, size=14, color=ACCENT, bold=True)
    add_text(s, Inches(7.2), y + Inches(0.15), Inches(5.3), Inches(1.4),
             desc, size=15, color=TEXT)
    y += Inches(1.7)

# ------- Slide 7: useRef -------
s = prs.slides.add_slide(blank)
add_bg(s)
add_accent_bar(s)
add_text(s, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
         "useRef — a mutable box that survives re-renders", size=32, bold=True, color=ACCENT)

add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.6),
         "const timersRef = useRef([]);", size=22, bold=True, color=ACCENT)

add_text(s, Inches(0.6), Inches(2.1), Inches(12), Inches(4.5),
         "• Holds the list of { type, id } timer handles created during a spin.\n"
         "• Not state — mutating it must NOT trigger a re-render.\n"
         "• Values persist across renders, unlike plain local variables.\n"
         "• scheduleSpinAnimation + staggered setTimeouts push IDs in.\n"
         "• clearTimers() cancels and resets the array when a new spin starts\n"
         "  or when the component unmounts (via useEffect cleanup).\n"
         "• Without the ref, cancelling mid-spin would be impossible —\n"
         "  the IDs would be lost between renders.",
         size=20, color=TEXT)

# ------- Slide 8: How they work together -------
s = prs.slides.add_slide(blank)
add_bg(s)
add_accent_bar(s)
add_text(s, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
         "How the four hooks cooperate in one spin", size=32, bold=True, color=ACCENT)

rows = [
    ("useState",  "Drives every visible change: spinning flag, indices, credits, message."),
    ("useRef",    "Stores raf/timeout IDs so we can cancel without re-rendering."),
    ("useMemo",   "Keeps the reel matrix + visible grid stable between state updates."),
    ("useEffect", "Guarantees cleanup when the component unmounts mid-animation."),
]
y = Inches(1.6)
for name, desc in rows:
    add_card(s, Inches(0.6), y, Inches(12.1), Inches(1.0))
    add_text(s, Inches(0.9), y + Inches(0.22), Inches(2.6), Inches(0.8),
             name, size=24, bold=True, color=ACCENT)
    add_text(s, Inches(3.6), y + Inches(0.25), Inches(9.0), Inches(0.8),
             desc, size=18, color=TEXT)
    y += Inches(1.2)

add_text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.6),
         "Together: smooth animation, accurate state, zero leaked timers.",
         size=18, color=SUB)

# ------- Slide 9: Thank you -------
s = prs.slides.add_slide(blank)
add_bg(s)
add_accent_bar(s)
add_text(s, Inches(0.6), Inches(2.7), Inches(12.1), Inches(1.4),
         "Thank you", size=72, bold=True, color=ACCENT,
         align=None)
add_text(s, Inches(0.6), Inches(4.2), Inches(12.1), Inches(0.8),
         "Questions?  ·  Live demo: npm run dev", size=24, color=TEXT)

out = '/Users/gustavboberg/develop/casionoslot/casino-slots.pptx'
prs.save(out)
print("Saved:", out)
